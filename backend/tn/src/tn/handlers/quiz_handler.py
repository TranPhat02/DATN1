"""
Quiz handler — Gemini AI RAG quiz generation + Manual quiz creation.
Uses MongoDB for quiz storage (flexible document schema).
Supports: multiple-choice (trắc nghiệm) and essay (tự luận) questions.
"""
import io
import json
import logging
import re
import traceback
import uuid
from typing import List, Optional
from datetime import datetime, timezone, timedelta

from fastapi import HTTPException
import google.generativeai as genai

from tn.handlers.drive_handler import get_content_for_quiz_ai, save_json_to_drive, _get_drive_service
from tn.handlers import memory_handler
from tn.config.database import settings, get_mongo_db, SessionLocal
from tn.models.diem_trac_nghiem import DiemTracNghiem

logger = logging.getLogger("quiz_handler")
logger.setLevel(logging.DEBUG)
if not logger.handlers:
    handler = logging.StreamHandler()
    handler.setFormatter(logging.Formatter("[QUIZ] %(levelname)s: %(message)s"))
    logger.addHandler(handler)


def _configure_gemini():
    """Configure Gemini AI with API key."""
    api_key = settings.GEMINI_API_KEY
    if not api_key:
        raise HTTPException(status_code=503, detail="GEMINI_API_KEY chưa được cấu hình")
    genai.configure(api_key=api_key)


def _extract_text_from_file(file_data: dict) -> str:
    """Extract readable text from a file based on its mimeType.
    Supports: PDF, DOCX, PPTX, and plain text files.
    """
    file_bytes = file_data["bytes"]
    mime = file_data.get("mimeType", "")
    name = file_data.get("name", "").lower()

    # ── PDF ──
    if mime == "application/pdf" or name.endswith(".pdf"):
        try:
            from PyPDF2 import PdfReader
            reader = PdfReader(io.BytesIO(file_bytes))
            pages_text = []
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    pages_text.append(text)
            result = "\n\n".join(pages_text)
            if result.strip():
                return result
        except Exception as e:
            print(f"[Quiz] PyPDF2 extraction failed: {e}")
        # Fallback: decode raw bytes (may include some junk but better than nothing)
        return file_bytes.decode("utf-8", errors="ignore")

    # ── DOCX ──
    if mime == "application/vnd.openxmlformats-officedocument.wordprocessingml.document" or name.endswith(".docx"):
        try:
            from docx import Document
            doc = Document(io.BytesIO(file_bytes))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            return "\n\n".join(paragraphs)
        except Exception as e:
            print(f"[Quiz] DOCX extraction failed: {e}")
            return file_bytes.decode("utf-8", errors="ignore")

    # ── PPTX ──
    if mime == "application/vnd.openxmlformats-officedocument.presentationml.presentation" or name.endswith(".pptx"):
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(file_bytes))
            slides_text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slides_text.append(shape.text)
            return "\n\n".join(slides_text)
        except Exception as e:
            print(f"[Quiz] PPTX extraction failed: {e}")
            return file_bytes.decode("utf-8", errors="ignore")

    # ── Plain text / other ──
    return file_bytes.decode("utf-8", errors="ignore")


def _parse_json_from_response(text: str) -> list:
    """Robustly extract JSON array from Gemini response text."""
    text = text.strip()
    # Remove markdown code fences
    if "```" in text:
        match = re.search(r"```(?:json)?\s*\n?(.*?)\n?\s*```", text, re.DOTALL)
        if match:
            text = match.group(1).strip()
    # Try direct parse
    try:
        return json.loads(text)
    except json.JSONDecodeError:
        pass
    # Try to find JSON array in text using manual index scan
    start = text.find('[')
    end = text.rfind(']')
    if start != -1 and end != -1 and end > start:
        try:
            return json.loads(text[start:end+1])
        except json.JSONDecodeError:
            pass
    raise ValueError(f"Không tìm thấy JSON hợp lệ trong phản hồi. Phản hồi gốc: {text[:300]}")


def _strip_option_prefix(text: str) -> str:
    """Strip leading A./B./C./D. (or A /, B /) prefixes from option text.
    Handles cases like 'A. Answer', 'A.Answer', 'A) Answer', 'A - Answer'.
    """
    import re
    cleaned = re.sub(r'^\s*[A-Da-d][.)\-:]\s*', '', text)
    return cleaned.strip() if cleaned.strip() else text.strip()


async def generate_quiz(file_id: str, ma_lop_mon: str, num_questions: int = 10, ma_gv: str = "") -> dict:
    """Generate quiz from a Drive file using Gemini RAG with Document Memory.
    
    Memory strategy (per file, namespaced by ma_lop_mon):
      - First call: AI reads 100% \u2192 saves content.md (summary) + quiz.md (pool)
      - Later calls: 70% from quiz.md pool + 30% regenerated from file chunk \u2192 updates pool
    """
    logger.info(f"=== QUIZ GENERATE START === file_id={file_id}, ma_lop_mon={ma_lop_mon}, num_questions={num_questions}")

    try:
        _configure_gemini()
        logger.info("[1/5] Gemini configured OK")
    except Exception as e:
        logger.error(f"[1/5] Gemini config FAILED: {e}\n{traceback.format_exc()}")
        raise

    # Download file content from Drive (with metadata for type detection)
    try:
        from tn.handlers.drive_handler import get_content_for_quiz_ai
        files_data = get_content_for_quiz_ai(file_id)
        logger.info(f"[2/5] Files downloaded: {len(files_data)} files")
    except Exception as e:
        logger.error(f"[2/5] Drive download FAILED: {e}\n{traceback.format_exc()}")
        raise HTTPException(status_code=400, detail=f"Kh\u00f4ng th\u1ec3 \u0111\u1ecdc t\u1ea3i d\u1eef li\u1ec7u t\u1eeb Google Drive: {str(e)}")

    try:
        file_text = ""
        # Collect file names for memory key
        file_names_for_key = []
        for f_data in files_data:
            text = _extract_text_from_file(f_data)
            fname = f_data.get('name', 'N/A')
            if text.strip():
                file_text += f"\n\n--- FILE: {fname} ---\n"
                file_text += text
                file_names_for_key.append(fname)
        logger.info(f"[3/5] Text extracted: {len(file_text)} chars, preview: {repr(file_text[:200])}")
    except Exception as e:
        logger.error(f"[3/5] Text extraction FAILED: {e}\n{traceback.format_exc()}")
        raise HTTPException(status_code=400, detail=f"Kh\u00f4ng th\u1ec3 tr\u00edch xu\u1ea5t text t\u1eeb ngu\u1ed3n: {str(e)}")

    if not file_text.strip():
        logger.error("[3/5] Extracted text is EMPTY")
        raise HTTPException(status_code=400, detail="File kh\u00f4ng c\u00f3 n\u1ed9i dung text \u0111\u1ec3 t\u1ea1o c\u00e2u h\u1ecfi")

    # Clean up text
    file_text = re.sub(r'\n{3,}', '\n\n', file_text)
    file_text = re.sub(r'[ \t]+', ' ', file_text)

    # Determine the memory key (use primary file name or joined names)
    memory_key = file_names_for_key[0] if len(file_names_for_key) == 1 else "_".join(file_names_for_key[:3])
    # Sanitize key (remove characters unsafe for folder names)
    import re as _re
    memory_key = _re.sub(r'[<>:"/\\|?*]', '_', memory_key)
    logger.info(f"[3/5] Memory key: '{memory_key}' for lop '{ma_lop_mon}'")

    # ── Step 4: Build/refresh content summary (memory) ──
    try:
        logger.info("[4/5] Building/loading content summary ...")
        # Use truncated text for summary (400k chars is enough for Gemini summary)
        summary_text = file_text[:400_000] if len(file_text) > 400_000 else file_text
        memory_handler.build_or_refresh_content_summary(
            ma_lop_mon=ma_lop_mon,
            file_name=memory_key,
            full_text=summary_text,
            force_rebuild=False,
        )
        logger.info("[4/5] Content summary OK")
    except Exception as e:
        # Non-blocking: if summary fails, continue with quiz generation
        logger.warning(f"[4/5] Content summary failed (non-blocking): {e}")

    # ── Step 5: Generate quiz using 70/30 memory strategy ──
    try:
        logger.info("[5/5] Generating quiz with memory ...")
        # Cap text to protect memory
        max_chars = 800_000
        if len(file_text) > max_chars:
            file_text = file_text[:max_chars]
            logger.info(f"[5/5] Text truncated to {max_chars} chars")

        questions_data = memory_handler.generate_quiz_with_memory(
            ma_lop_mon=ma_lop_mon,
            file_name=memory_key,
            full_text=file_text,
            num_questions=num_questions,
        )
        logger.info(f"[5/5] Parsed {len(questions_data)} questions OK")
    except HTTPException:
        raise
    except Exception as e:
        error_msg = str(e)
        logger.error(f"[5/5] Quiz generation FAILED: {type(e).__name__}: {error_msg}\n{traceback.format_exc()}")

        if "429" in error_msg or "ResourceExhausted" in error_msg or "quota" in error_msg.lower():
            raise HTTPException(status_code=429, detail="Gemini API: \u0110\u00e3 v\u01b0\u1ee3t qu\u00e1 gi\u1edbi h\u1ea1n h\u1ea1n m\u1ee9c (Quota) mi\u1ec5n ph\u00ed. H\u00e3y th\u1eed l\u1ea1i sau!")

        raise HTTPException(status_code=500, detail=f"Gemini kh\u00f4ng t\u1ea1o \u0111\u01b0\u1ee3c c\u00e2u h\u1ecfi: {error_msg}")

    if not questions_data:
        raise HTTPException(status_code=500, detail="AI kh\u00f4ng sinh \u0111\u01b0\u1ee3c c\u00e2u h\u1ecfi n\u00e0o t\u1eeb t\u00e0i li\u1ec7u n\u00e0y")

    # Build quiz document
    quiz_id = str(uuid.uuid4())[:8]
    quiz = {
        "id": quiz_id,
        "maLopMon": ma_lop_mon,
        "maGV": ma_gv,
        "title": f"Tr\u1eafc nghi\u1ec7m AI - {ma_lop_mon} - {datetime.now().strftime('%d/%m/%Y')}",
        "type": "ai",
        "questions": [
            {
                "id": str(uuid.uuid4())[:8],
                "type": "multiple_choice",
                "question": q["question"],
                "options": q["options"],
                "correctAnswer": q["correctAnswer"],
                "explanation": q.get("explanation", ""),
            }
            for q in questions_data
        ],
        "duration": max(num_questions * 2, 10),  # 2 min per question, min 10
        "fileId": file_id,
        "memoryKey": memory_key,
        "lockType": 1,
        "lockUntil": "",
        "createdAt": datetime.now().isoformat(),
    }

    logger.info(f"=== QUIZ GENERATE DONE === quiz_id={quiz_id} (Not Auto-Saved)")
    return quiz


async def create_manual_quiz(ma_lop_mon: str, title: str, duration: int, questions: list, ma_gv: str = "", lock_type: int = 1, lock_until: str = "") -> dict:
    """Create a manual quiz with mixed question types (multiple_choice / essay)."""
    quiz_id = str(uuid.uuid4())[:8]

    processed_questions = []
    for q in questions:
        q_type = q.get("type", "multiple_choice")
        question_doc = {
            "id": q.get("id") or str(uuid.uuid4())[:8],
            "type": q_type,
            "question": q["question"],
        }
        if "imageUrl" in q and q["imageUrl"]:
            question_doc["imageUrl"] = q["imageUrl"]
            
        if q_type == "multiple_choice":
            question_doc["options"] = q.get("options", [])
            question_doc["correctAnswer"] = q.get("correctAnswer", 0)
            question_doc["explanation"] = q.get("explanation", "")
        elif q_type == "essay":
            question_doc["sampleAnswer"] = q.get("sampleAnswer", "")
            question_doc["maxScore"] = q.get("maxScore", 10.0)
        processed_questions.append(question_doc)

    quiz = {
        "id": quiz_id,
        "maLopMon": ma_lop_mon,
        "maGV": ma_gv,
        "title": title or f"B\u00e0i ki\u1ec3m tra - {datetime.now().strftime('%d/%m/%Y')}",
        "type": "manual",
        "questions": processed_questions,
        "duration": duration or 30,
        "lockType": 1,        # Always locked on creation — teacher must unlock manually
        "lockUntil": "",
        "createdAt": datetime.now().isoformat(),
    }

    db = get_mongo_db()
    if db is not None:
        await db.quizzes.insert_one(quiz.copy())

    # Save quiz to Google Drive at /quiz
    try:
        drive_quiz = quiz.copy()
        drive_quiz.pop("_id", None)
        save_json_to_drive(ma_lop_mon, ["quiz"], f"quiz_{quiz_id}.json", drive_quiz)
        logger.info(f"Manual quiz saved to Drive: /quiz/quiz_{quiz_id}.json")
    except Exception as e:
        logger.warning(f"Failed to save manual quiz to Drive (non-blocking): {e}")

    return quiz


async def edit_quiz(quiz_id: str, ma_lop_mon: str, title: str, duration: int, questions: list, ma_gv: str = "", lock_type: int = 0, lock_until: str = "") -> dict:
    """Edit an existing quiz in MongoDB."""
    db = get_mongo_db()
    if db is None:
        raise HTTPException(status_code=500, detail="MongoDB chưa sẵn sàng")
        
    old_quiz = await db.quizzes.find_one({"id": quiz_id})
    if not old_quiz:
        raise HTTPException(status_code=404, detail="Không tìm thấy quiz")

    processed_questions = []
    for q in questions:
        q_type = q.get("type", "multiple_choice")
        question_doc = {
            "id": q.get("id") or str(uuid.uuid4())[:8],
            "type": q_type,
            "question": q["question"],
        }
        if "imageUrl" in q and q["imageUrl"]:
            question_doc["imageUrl"] = q["imageUrl"]
            
        if q_type == "multiple_choice":
            question_doc["options"] = q.get("options", [])
            question_doc["correctAnswer"] = q.get("correctAnswer", 0)
            question_doc["explanation"] = q.get("explanation", "")
        elif q_type == "essay":
            question_doc["sampleAnswer"] = q.get("sampleAnswer", "")
            question_doc["maxScore"] = q.get("maxScore", 10.0)
        processed_questions.append(question_doc)

    updates = {
        "title": title or old_quiz.get("title", ""),
        "type": "manual",  # If edited, consider it manual from now on
        "questions": processed_questions,
        "duration": duration or 30,
        "lockType": lock_type,
        "lockUntil": lock_until,
        "updatedAt": datetime.now().isoformat(),
    }

    await db.quizzes.update_one({"id": quiz_id}, {"$set": updates})
    
    # Reload and return
    updated_quiz = await db.quizzes.find_one({"id": quiz_id}, {"_id": 0})
    
    # Save to Drive again if logic applies
    try:
        drive_quiz = updated_quiz.copy()
        save_json_to_drive(ma_lop_mon, ["quiz"], f"quiz_{quiz_id}.json", drive_quiz)
    except Exception as e:
        logger.warning(f"Failed to save edited quiz to Drive: {e}")

    return updated_quiz


async def list_quizzes_by_lop_mon(ma_lop_mon: str) -> List[dict]:
    """Get all quizzes for a class from MongoDB."""
    db = get_mongo_db()
    if db is None:
        return []
    cursor = db.quizzes.find({"maLopMon": ma_lop_mon}, {"_id": 0})
    quizzes = []
    now = datetime.now(timezone(timedelta(hours=7)))
    async for doc in cursor:
        # Hide correct answers and explanations for student view
        for q in doc.get("questions", []):
            q.pop("correctAnswer", None)
            q.pop("explanation", None)
            q.pop("sampleAnswer", None)
        # Compute lock status
        lock_type = doc.get("lockType", 0)
        lock_until = doc.get("lockUntil", "")
        if lock_type == 1:
            doc["isLocked"] = True
        elif lock_type == 2 and lock_until:
            try:
                deadline = datetime.fromisoformat(lock_until)
                if deadline.tzinfo is None:
                    deadline = deadline.replace(tzinfo=timezone(timedelta(hours=7)))
                doc["isLocked"] = now > deadline
            except Exception:
                doc["isLocked"] = False
        else:
            doc["isLocked"] = False
        quizzes.append(doc)
    return quizzes


async def get_quiz_by_id(quiz_id: str) -> dict:
    """Get a single quiz (without correct answers for students)."""
    db = get_mongo_db()
    if db is None:
        raise HTTPException(status_code=404, detail="Quiz không tìm thấy")
    quiz = await db.quizzes.find_one({"id": quiz_id}, {"_id": 0})
    if not quiz:
        raise HTTPException(status_code=404, detail="Quiz không tìm thấy")
    
    # Check lock status before returning quiz data
    lock_type = quiz.get("lockType", 0)
    lock_until = quiz.get("lockUntil", "")
    now = datetime.now(timezone(timedelta(hours=7)))
    is_locked = False
    if lock_type == 1:
        is_locked = True
    elif lock_type == 2 and lock_until:
        try:
            deadline = datetime.fromisoformat(lock_until)
            if deadline.tzinfo is None:
                deadline = deadline.replace(tzinfo=timezone(timedelta(hours=7)))
            is_locked = now > deadline
        except Exception:
            pass
    if is_locked:
        raise HTTPException(status_code=403, detail="Bài trắc nghiệm đang bị khoá")
    
    # Remove answers + explanation
    student_quiz = quiz.copy()
    for q in student_quiz.get("questions", []):
        q.pop("correctAnswer", None)
        q.pop("explanation", None)
        q.pop("sampleAnswer", None)
    return student_quiz


async def update_quiz_lock(quiz_id: str, lock_type: int, lock_until: str) -> dict:
    """Update lock settings for a quiz."""
    db = get_mongo_db()
    if db is None:
        raise HTTPException(status_code=500, detail="MongoDB không khả dụng")
    result = await db.quizzes.update_one(
        {"id": quiz_id},
        {"$set": {"lockType": lock_type, "lockUntil": lock_until}}
    )
    if result.matched_count == 0:
        raise HTTPException(status_code=404, detail="Quiz không tìm thấy")
    return {"message": "Cập nhật khoá thành công", "lockType": lock_type, "lockUntil": lock_until}


async def submit_quiz(quiz_id: str, ma_sv: str, answers: dict, tab_switch_count: int, time_spent: int = 0) -> dict:
    """Grade a quiz submission and save score to MongoDB."""
    db = get_mongo_db()
    if db is None:
        raise HTTPException(status_code=500, detail="MongoDB chưa sẵn sàng")

    # Get quiz with answers
    quiz = await db.quizzes.find_one({"id": quiz_id}, {"_id": 0})
    if not quiz:
        raise HTTPException(status_code=404, detail="Quiz không tìm thấy")

    # Verify quiz is not locked at the time of submission
    lock_type = quiz.get("lockType", 0)
    lock_until = quiz.get("lockUntil", "")
    now = datetime.now(timezone(timedelta(hours=7)))
    if lock_type == 1:
        raise HTTPException(status_code=403, detail="Bài thi đã bị khoá cứng, không thể nộp bài")
    elif lock_type == 2 and lock_until:
        try:
            deadline = datetime.fromisoformat(lock_until)
            if deadline.tzinfo is None:
                deadline = deadline.replace(tzinfo=timezone(timedelta(hours=7)))
            # Add 1 minute grace period for latency during frontend auto-submit
            if now > (deadline + timedelta(minutes=1)):
                raise HTTPException(status_code=403, detail="Bài thi đã hết hạn, không thể nộp bài")
        except Exception:
            pass

    # Grade multiple choice questions
    correct = 0
    mc_total = 0
    mc_score = 0.0
    essay_answers = []
    
    total_q_len = len(quiz.get("questions", []))
    default_weight = 10.0 / total_q_len if total_q_len > 0 else 0.0

    for q in quiz["questions"]:
        weight = float(q.get("maxScore", default_weight))
        if q.get("type", "multiple_choice") == "multiple_choice":
            mc_total += 1
            student_answer = answers.get(q["id"])
            if student_answer is not None and student_answer == q.get("correctAnswer"):
                mc_score += weight
                correct += 1
        elif q.get("type") == "essay":
            essay_answers.append({
                "questionId": q["id"],
                "answer": answers.get(q["id"], ""),
            })

    total = mc_total + len(essay_answers)
    
    # Calculate initial score based on MC answers
    initial_score = round(mc_score, 2)
    
    # Apply tab switch penalties
    tab_penalty = 0.0
    status_flag = ""
    if tab_switch_count > 5:
        initial_score = 0.0
        status_flag = "Cảnh cáo: 0 điểm do chuyển tab quá 5 lần"
    else:
        if tab_switch_count > 4:
            tab_penalty = 1.0
            status_flag = "Cảnh cáo Lần 2: Trừ 1.0 điểm do chuyển tab quá 4 lần"
        elif tab_switch_count > 2:
            tab_penalty = 0.5
            status_flag = "Cảnh cáo Lần 1: Trừ 0.5 điểm do chuyển tab quá 2 lần"

    score = max(0.0, initial_score - tab_penalty)

    # Check if a prior submission exists (retake logic)
    prior_submission = await db.quiz_submissions.find_one({"quizId": quiz_id, "maSV": ma_sv})
    is_retake = prior_submission is not None

    has_essays = len(essay_answers) > 0
    status = "pending" if has_essays else "graded"

    if not is_retake:
        # Only save to DB and Drive for the first attempt
        submission = {
            "quizId": quiz_id,
            "maSV": ma_sv,
            "maLopMon": quiz.get("maLopMon"),
            "soCauDung": correct,
            "tongSoCau": total,
            "mcTotal": mc_total,
            "score": round(score, 1),
            "answers": answers,
            "essayAnswers": essay_answers,
            "tabSwitchCount": tab_switch_count,
            "timeSpent": time_spent,
            "status": status,
            "statusFlag": status_flag,
            "submittedAt": datetime.now(timezone(timedelta(hours=7))).isoformat(),
        }
        await db.quiz_submissions.insert_one(submission.copy())

        # If fully graded (no essays), save directly to MySQL DiemTracNghiem
        if status == "graded":
            try:
                db_sql = SessionLocal()
                # Use a specific composite ID for MaTN, or just a new UUID. We'll use quiz_id + ma_sv to avoid conflict? 
                # Actually MaTN is primary key, let's just make it shorter or use the quiz_id if 1 quiz = 1 score.
                # Since a student can take multiple quizzes in one lop_mon, DiemTracNghiem is usually 1 row per quiz per student.
                # MaTN in DiemTracNghiem is String(20). We can use f"{quiz_id}_{ma_sv}"[:20].
                ma_tn_val = f"{quiz_id}_{ma_sv}"
                vn_tz = timezone(timedelta(hours=7))
                existing = db_sql.query(DiemTracNghiem).filter(DiemTracNghiem.MaTN == ma_tn_val).first()
                if existing:
                    existing.SoCauDung = correct
                    existing.TongSoCau = total
                    existing.ThoiGianLam = time_spent
                    existing.ThoiGianNop = datetime.now(vn_tz)
                    existing.SoLanViPham = tab_switch_count
                    existing.FileID = quiz.get("title")
                else:
                    new_diem = DiemTracNghiem(
                        MaTN=ma_tn_val,
                        MaSV=ma_sv,
                        MaLopMon=quiz.get("maLopMon"),
                        SoCauDung=correct,
                        TongSoCau=total,
                        ThoiGianLam=time_spent,
                        ThoiGianNop=datetime.now(vn_tz),
                        SoLanViPham=tab_switch_count,
                        FileID=quiz.get("title"),
                    )
                    db_sql.add(new_diem)
                db_sql.commit()
            except Exception as e:
                logger.error(f"Failed to save score to MySQL: {e}")
            finally:
                db_sql.close()
    
        # Save result to Google Drive at /ket_qua/{quiz_id}/{ma_sv}.json
        ma_lop_mon = quiz.get("maLopMon", "")
        if ma_lop_mon:
            try:
                drive_result = submission.copy()
                drive_result.pop("_id", None)
                save_json_to_drive(
                    ma_lop_mon,
                    ["ket_qua", quiz_id],
                    f"{ma_sv}.json",
                    drive_result,
                )
                logger.info(f"Quiz result saved to Drive: /ket_qua/{quiz_id}/{ma_sv}.json")
            except Exception as e:
                logger.warning(f"Failed to save quiz result to Drive (non-blocking): {e}")

    return {
        "maTN": quiz_id,
        "soCauDung": correct,
        "tongSoCau": total,
        "score": round(score, 1),
        "isRetake": is_retake,
    }




async def delete_quiz(quiz_id: str) -> dict:
    """Delete a quiz from MongoDB."""
    db = get_mongo_db()
    if db is not None:
        await db.quizzes.delete_one({"id": quiz_id})
    return {"message": f"Quiz '{quiz_id}' đã xóa"}


async def list_pending_submissions(ma_lop_mon: str) -> List[dict]:
    """Get all submissions that have pending essays to grade for a class."""
    db = get_mongo_db()
    if db is None:
        raise HTTPException(status_code=500, detail="MongoDB chưa sẵn sàng")
    
    submissions = await db.quiz_submissions.find({"maLopMon": ma_lop_mon, "status": "pending"}).to_list(length=None)
    for sub in submissions:
        sub.pop("_id", None)
    return submissions


async def grade_submission(quiz_id: str, ma_sv: str, essay_grades: dict) -> dict:
    """Grade an essay submission and finalize the score in MySQL DiemTracNghiem."""
    db = get_mongo_db()
    if db is None:
        raise HTTPException(status_code=500, detail="MongoDB chưa sẵn sàng")

    submission = await db.quiz_submissions.find_one({"quizId": quiz_id, "maSV": ma_sv, "status": "pending"})
    if not submission:
        raise HTTPException(status_code=404, detail="Không tìm thấy bài làm cần chấm hoặc bài đã chấm")
    
    quiz = await db.quizzes.find_one({"id": quiz_id})
    if not quiz:
        raise HTTPException(status_code=404, detail="Không tìm thấy quiz liên quan")
        
    mc_score = submission.get("score", 0.0) # MC score natively out of 10
    
    # Calculate earned essay points
    earned_essay_points = sum(essay_grades.values())
    
    # Final score out of 10
    final_score = mc_score + earned_essay_points
    if final_score > 10.0: final_score = 10.0
    
    total_q = submission.get("tongSoCau", 1)
    
    # Map back to real total questions so it looks accurate on the UI.
    fake_tong_so = total_q
    fake_so_dung = round((final_score / 10) * total_q)
    
    # Mark as graded in Mongo
    await db.quiz_submissions.update_one(
        {"quizId": quiz_id, "maSV": ma_sv},
        {"$set": {"status": "graded", "finalScore": round(final_score, 1), "essayGrades": essay_grades}}
    )
    
    # Save to MySQL table DiemTracNghiem
    try:
        db_sql = SessionLocal()
        ma_tn_val = f"{quiz_id}_{ma_sv}"
        existing = db_sql.query(DiemTracNghiem).filter(DiemTracNghiem.MaTN == ma_tn_val).first()
        if not existing:
            new_diem = DiemTracNghiem(
                MaTN=ma_tn_val,
                MaSV=ma_sv,
                MaLopMon=submission.get("maLopMon"),
                SoCauDung=fake_so_dung,
                TongSoCau=fake_tong_so,
                SoLanViPham=submission.get("tabSwitchCount", 0),
                FileID=quiz.get("title"),
            )
            db_sql.add(new_diem)
        else:
            # Overwrite if exists, although it shouldn't for initial grade
            existing.SoCauDung = fake_so_dung
            existing.TongSoCau = fake_tong_so
            existing.SoLanViPham = submission.get("tabSwitchCount", 0)
            existing.FileID = quiz.get("title")
        db_sql.commit()
    except Exception as e:
        logger.error(f"Failed to save graded score to MySQL: {e}")
    finally:
        db_sql.close()
        
    return {"message": "Đã chấm điểm thành công", "finalScore": round(final_score, 1)}


async def get_submission_detail(quiz_id: str, ma_sv: str, caller_username: str = "", caller_role: str = "") -> dict:
    """Get full details of a submission including the original quiz questions and all answers."""
    db = get_mongo_db()
    if db is None:
        raise HTTPException(status_code=500, detail="MongoDB chưa sẵn sàng")

    # Students can only view their own submissions
    if caller_role == "student" and caller_username and ma_sv.upper() != caller_username.upper():
        raise HTTPException(status_code=403, detail="Bạn chỉ có thể xem bài làm của chính mình")

    submission = await db.quiz_submissions.find_one({"quizId": quiz_id, "maSV": ma_sv}, {"_id": 0})
    if not submission:
        raise HTTPException(status_code=404, detail="Không tìm thấy bài làm")
        
    quiz = await db.quizzes.find_one({"id": quiz_id}, {"_id": 0})
    if not quiz:
        raise HTTPException(status_code=404, detail="Không tìm thấy thông tin bài trắc nghiệm")
        
    return {
        "submission": submission,
        "quiz": quiz
    }


async def list_all_submissions(ma_lop_mon: str) -> List[dict]:
    """Get all submissions for an entire class."""
    db = get_mongo_db()
    if db is None:
        raise HTTPException(status_code=500, detail="MongoDB chưa sẵn sàng")
        
    cursor = db.quiz_submissions.find({"maLopMon": ma_lop_mon}, {"_id": 0})
    submissions = await cursor.to_list(length=None)
    return submissions


async def get_my_submissions(ma_lop_mon: str, ma_sv: str) -> List[dict]:
    """Return a list of {quizId, status} for quizzes that a student has submitted."""
    db = get_mongo_db()
    if db is None:
        return []
    cursor = db.quiz_submissions.find(
        {"maLopMon": ma_lop_mon, "maSV": ma_sv},
        {"quizId": 1, "status": 1, "_id": 0}
    )
    docs = await cursor.to_list(length=None)
    return [{"quizId": doc["quizId"], "status": doc.get("status", "graded")} for doc in docs if "quizId" in doc]


async def list_submissions_grouped_by_quiz(ma_lop_mon: str) -> List[dict]:
    """Get all submissions for a class, grouped by quiz, enriched with quiz title and TenMH.

    Returns a list of quiz groups:
    [
      {
        "quizId": str,
        "quizTitle": str,
        "tenMH": str,
        "maLopMon": str,
        "totalSubmissions": int,
        "submissions": [
          {
            "maSV": str,
            "tenSV": str (resolved),
            "soCauDung": int,
            "tongSoCau": int,
            "diem": float,
            "thoiGianLam": int,
            "thoiGianNop": str,
            "soLanViPham": int,
            "status": str,
          },
          ...
        ]
      }
    ]
    """
    db = get_mongo_db()
    if db is None:
        raise HTTPException(status_code=500, detail="MongoDB chưa sẵn sàng")

    # 1. Load TenMH from MySQL (LopMonHoc → MonHoc)
    ten_mh = ma_lop_mon  # fallback
    try:
        db_sql = SessionLocal()
        from tn.models.lop_mon_hoc import LopMonHoc
        from tn.models.mon_hoc import MonHoc
        lmh = db_sql.query(LopMonHoc).filter(LopMonHoc.MaLopMon == ma_lop_mon).first()
        if lmh and lmh.MaMH:
            mh = db_sql.query(MonHoc).filter(MonHoc.MaMH == lmh.MaMH).first()
            if mh:
                ten_mh = mh.TenMH or lmh.MaMH
        db_sql.close()
    except Exception as e:
        logger.warning(f"Could not resolve TenMH for {ma_lop_mon}: {e}")

    # 2. Load all submissions for this class from MongoDB
    cursor = db.quiz_submissions.find({"maLopMon": ma_lop_mon}, {"_id": 0})
    all_submissions = await cursor.to_list(length=None)

    # 3. Load all quizzes for this class (for their titles)
    quiz_cursor = db.quizzes.find({"maLopMon": ma_lop_mon}, {"_id": 0, "id": 1, "title": 1})
    quizzes_raw = await quiz_cursor.to_list(length=None)
    quiz_title_map: dict[str, str] = {q["id"]: q.get("title", q["id"]) for q in quizzes_raw}

    # 4. Resolve student names from MySQL
    try:
        db_sql = SessionLocal()
        from tn.models.sinh_vien import SinhVien
        students = db_sql.query(SinhVien).all()
        sv_name_map: dict[str, str] = {sv.MaSV: sv.TenSV or sv.MaSV for sv in students}
        db_sql.close()
    except Exception as e:
        logger.warning(f"Could not load student names: {e}")
        sv_name_map = {}

    # 5. Group by quizId
    groups: dict[str, dict] = {}
    for sub in all_submissions:
        quiz_id = sub.get("quizId", "")
        if not quiz_id:
            continue

        if quiz_id not in groups:
            groups[quiz_id] = {
                "quizId": quiz_id,
                "quizTitle": quiz_title_map.get(quiz_id, f"Bài TN {quiz_id}"),
                "tenMH": ten_mh,
                "maLopMon": ma_lop_mon,
                "submissions": [],
            }

        ma_sv = sub.get("maSV", "")
        so_cau_dung = sub.get("soCauDung", 0) or 0
        tong_so_cau = sub.get("tongSoCau", 1) or 1
        score = round(sub.get("score", (so_cau_dung / tong_so_cau) * 10), 1) if tong_so_cau > 0 else 0.0

        # Use finalScore if essay was graded
        if sub.get("finalScore") is not None:
            score = round(sub["finalScore"], 1)

        groups[quiz_id]["submissions"].append({
            "maSV": ma_sv,
            "tenSV": sv_name_map.get(ma_sv, ma_sv),
            "soCauDung": so_cau_dung,
            "tongSoCau": tong_so_cau,
            "diem": score,
            "thoiGianLam": sub.get("timeSpent") or sub.get("thoiGianLam") or 0,
            "thoiGianNop": sub.get("submittedAt", ""),
            "soLanViPham": sub.get("tabSwitchCount", 0) or 0,
            "status": sub.get("status", "graded"),
        })

    # 6. Sort groups by quizId descending (newest first), compute totalSubmissions
    result = sorted(groups.values(), key=lambda g: g["quizId"], reverse=True)
    for g in result:
        g["totalSubmissions"] = len(g["submissions"])

    return result
